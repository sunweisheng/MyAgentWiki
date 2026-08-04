from __future__ import annotations

from pathlib import Path
from unittest import mock

from docx import Document
from openpyxl import Workbook
from pptx import Presentation
from pptx.util import Inches

from myagentwiki.app_services import runtime_services
from myagentwiki.app_services.document_conversion import convert_local_document
from myagentwiki.cli import infer_source_type


def test_markitdown_converts_common_document_formats(tmp_path: Path) -> None:
    docx_path = tmp_path / "sample.docx"
    document = Document()
    document.add_heading("Word 标题", level=1)
    document.add_paragraph("Word 正文需要转换成 Markdown。")
    document.save(docx_path)

    xlsx_path = tmp_path / "sample.xlsx"
    workbook = Workbook()
    sheet = workbook.active
    sheet.title = "数据"
    sheet.append(["名称", "数量"])
    sheet.append(["样例", 2])
    workbook.save(xlsx_path)

    pptx_path = tmp_path / "sample.pptx"
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[5])
    slide.shapes.title.text = "演示标题"
    text_box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(5), Inches(1))
    text_box.text_frame.text = "演示正文需要进入 Markdown。"
    presentation.save(pptx_path)

    html_path = tmp_path / "sample.html"
    html_path.write_text("<h1>网页标题</h1><p>网页正文需要转换成 Markdown。</p>", encoding="utf-8")

    expectations = {
        docx_path: "Word 正文",
        xlsx_path: "样例",
        pptx_path: "演示正文",
        html_path: "网页正文",
    }
    for path, expected_text in expectations.items():
        markdown, metadata = convert_local_document(path)
        assert expected_text in markdown
        assert metadata["extraction_method"] == "markitdown"
        assert metadata["location_map"]["converter"] == "microsoft/markitdown"


def test_document_dispatch_uses_markitdown_as_primary_converter(tmp_path: Path) -> None:
    raw_path = tmp_path / "source.html"
    raw_path.write_text("<h1>统一转换</h1><p>由 MarkItDown 负责文档转 Markdown。</p>", encoding="utf-8")

    markdown, metadata = runtime_services.convert_source_to_normalized_markdown(raw_path, "html")

    assert "由 MarkItDown 负责文档转 Markdown" in markdown
    assert metadata["extraction_method"] == "markitdown"
    assert metadata["warnings"] == []


def test_document_dispatch_records_markitdown_failure_before_legacy_fallback(tmp_path: Path) -> None:
    raw_path = tmp_path / "legacy.doc"
    raw_path.write_bytes(
        b"\xd0\xcf\x11\xe0\xa1\xb1\x1a\xe1"
        + b"\x00\x01BinaryPrefix"
        + b"Project Plan 2026"
    )

    with mock.patch.object(
        runtime_services.document_conversion,
        "convert_local_document",
        side_effect=RuntimeError("mock MarkItDown failure"),
    ):
        markdown, metadata = runtime_services.convert_source_to_normalized_markdown(raw_path, "doc")

    assert "Project Plan 2026" in markdown
    assert metadata["extraction_method"] == "markitdown_failed+python_only"
    assert "markitdown_conversion_failed:RuntimeError" in metadata["warnings"]
    assert metadata["location_map"]["markitdown"]["status"] == "failed"
    assert metadata["location_map"]["markitdown"]["converter"] == "microsoft/markitdown"
    assert metadata["location_map"]["markitdown"]["converter_version"]


def test_normalized_record_currentness_checks_pipeline_and_markitdown_versions() -> None:
    record = {
        "normalizer_version": runtime_services.NORMALIZER_VERSION,
        "extraction_method": "markitdown",
        "location_map": {
            "converter": "microsoft/markitdown",
            "converter_version": runtime_services.document_conversion.get_markitdown_version(),
        },
    }

    assert runtime_services.normalized_record_is_current(record) is True

    record["normalizer_version"] = "normalize_v1"
    assert runtime_services.normalized_record_is_current(record) is False

    record["normalizer_version"] = runtime_services.NORMALIZER_VERSION
    with mock.patch.object(
        runtime_services.document_conversion,
        "get_markitdown_version",
        return_value="new-version",
    ):
        assert runtime_services.normalized_record_is_current(record) is False


def test_infer_source_type_covers_markitdown_document_families(tmp_path: Path) -> None:
    expected_types = {
        "slides.pptx": "presentation",
        "page.html": "html",
        "data.json": "structured_text",
        "book.epub": "ebook",
        "notes.ipynb": "notebook",
        "mail.msg": "email",
    }

    for filename, expected_type in expected_types.items():
        assert infer_source_type(tmp_path / filename) == expected_type


def test_runtime_services_legacy_doc_converter_extracts_visible_snippets(tmp_path) -> None:
    raw_path = tmp_path / "legacy.doc"
    payload = (
        b"\xd0\xcf\x11\xe0\xa1\xb1\x1a\xe1"
        + b"\x00\x01BinaryPrefix"
        + b"Project Plan 2026"
        + b"\x00\x02"
        + "知识声明层需要可追踪".encode("utf-16le")
        + b"\x00\x03Tail"
    )
    raw_path.write_bytes(payload)

    markdown, metadata = runtime_services.convert_legacy_doc_to_markdown(raw_path)

    assert "# legacy" in markdown
    assert metadata["extraction_quality"] == "partial"
    assert "legacy_doc_binary_fallback" in metadata["warnings"]


def test_runtime_services_image_converter_without_tesseract_keeps_metadata_only(tmp_path) -> None:
    raw_path = tmp_path / "diagram.png"
    raw_path.write_bytes(
        b"\x89PNG\r\n\x1a\n"
        + b"\x00\x00\x00\x0dIHDR"
        + (32).to_bytes(4, "big")
        + (16).to_bytes(4, "big")
        + b"\x08\x02\x00\x00\x00"
        + b"\x00\x00\x00\x00"
    )

    with mock.patch("myagentwiki.app_services.runtime_services.command_exists", return_value=False):
        markdown, metadata = runtime_services.convert_image_to_markdown(raw_path)

    assert "# diagram" in markdown
    assert metadata["extraction_quality"] == "partial"
    assert "tesseract_missing" in metadata["warnings"]
